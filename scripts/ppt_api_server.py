from flask import Flask, request, send_file, jsonify
import google.generativeai as genai
import json
from pptx import Presentation
from pptx.util import Inches, Pt
import tempfile
import os

API_KEY = "AIzaSyD0x2jEZp73JOWp1rjBcw7jBYVNHmUVkys"
genai.configure(api_key=API_KEY)

app = Flask(__name__)

def fetch_from_gemini(user_input):
    generation_config = {
        "temperature": 1,
        "top_p": 0.95,
        "top_k": 64,
        "max_output_tokens": 2500,
    }
    model = genai.GenerativeModel(
        model_name="gemini-1.5-flash",
        generation_config=generation_config,
    )
    chat_session = model.start_chat(history=[{"role": "user", "parts": [user_input]}])
    response = chat_session.send_message(user_input)
    return response.text

def extract_slides_json(result):
    import re
    match = re.search(r"```json[\s\S]*?(\{[\s\S]*'slides'[\s\S]*\})[\s\S]*?```", result)
    if match:
        json_str = match.group(1)
    else:
        curly_match = re.search(r"\{[\s\S]*'slides'[\s\S]*\}", result)
        json_str = curly_match.group(0) if curly_match else result
    json_str = json_str.replace("'", '"')
    try:
        slides_obj = json.loads(json_str)
        return slides_obj.get('slides', [])
    except Exception as e:
        print("Failed to parse slides JSON:", e)
        print("Raw output:\n", result)
        return []

def create_pptx(slides):
    prs = Presentation()
    for slide_data in slides:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        title.text = slide_data.get('title', 'Slide')
        content = slide_data.get('content', '')
        textbox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(4))
        tf = textbox.text_frame
        for line in content.split('\n'):
            p = tf.add_paragraph()
            p.text = line.strip()
            p.font.size = Pt(18)
    temp_file = tempfile.NamedTemporaryFile(delete=False, suffix='.pptx')
    prs.save(temp_file.name)
    temp_file.close()
    return temp_file.name

@app.route('/generate-ppt', methods=['POST'])
def generate_ppt():
    data = request.get_json()
    project_idea = data.get('project_idea')
    if not project_idea:
        return jsonify({'error': 'Missing project_idea'}), 400
    prompt = f"""
You are an expert product manager and startup consultant. A user has provided a project idea.\nYour task is to generate the content for a 6-slide presentation deck based on this idea.\nRespond ONLY with a valid JSON object containing a single key 'slides', which is an array of 6 slide objects.\nEach slide object must have a 'title' and 'content' property. The content for each slide should be concise, professional, and formatted in bullet points.\nThe required slide titles are exactly:\n1. Problem Definition\n2. Technical Approach\n3. Impacts and Benefits\n4. Feasibility and Viability\n5. Research Review\n6. Conclusion & Next Steps\nHere is the user's project idea:\n---\n{project_idea}\n---\nRespond ONLY with valid JSON, no explanation, no markdown, no extra text, no code block, no comments, no formatting. Output only the JSON object."""
    result = fetch_from_gemini(prompt)
    slides = extract_slides_json(result)
    if not slides:
        return jsonify({'error': 'Failed to generate slides'}), 500
    pptx_path = create_pptx(slides)
    response = send_file(pptx_path, as_attachment=True, download_name='pitch_deck.pptx')
    # Clean up temp file after sending
    os.remove(pptx_path)
    return response

if __name__ == '__main__':
    app.run(port=5000, debug=True)
